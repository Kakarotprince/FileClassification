# utils.py
import os
import fitz  # PyMuPDF
#import textract
import pandas as pd
from bs4 import BeautifulSoup
from docx import Document
from pptx import Presentation
from sklearn.base import BaseEstimator, ClassifierMixin
import re
from PIL import Image
import pytesseract
import nltk
import numpy as np


# Ensure NLTK downloads
nltk.download('punkt', quiet=True)

def clean_text(text):
    text = text.lower()
    text = re.sub(r'[^a-zA-Z\s]', '', text)
    return text

def extract_text_from_file(file_path):
    _, ext = os.path.splitext(file_path)
    ext = ext.lower()

    try:
        if ext == '.pdf':
            doc = fitz.open(file_path)
            text = ""
            for page in doc:
                page_text = page.get_text()
                if not page_text.strip():
                    pix = page.get_pixmap(dpi=300)
                    img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                    page_text = pytesseract.image_to_string(img)
                text += page_text + "\n"
            return text

        elif ext == '.docx':
            doc = Document(file_path)
            return "\n".join(para.text for para in doc.paragraphs)

        elif ext in ['.xls', '.xlsx']:
            df = pd.read_excel(file_path, dtype=str)
            return df.apply(" ".join, axis=1).str.cat(sep="\n")
        
        elif ext == '.pptx':
            if not os.path.exists(file_path):
                return f"File not found: {file_path}"
            try:
                prs = Presentation(file_path)
                text = ""
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text += shape.text + "\n"
                return text
            except Exception as e:
                return f"Could not extract text from PPTX: {e}"


        elif ext == '.txt':
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read()

        elif ext in ['.html', '.htm']:
            with open(file_path, 'r', encoding='utf-8') as f:
                soup = BeautifulSoup(f, 'html.parser')
                return soup.get_text()

        elif ext in ['.jpg', '.jpeg', '.png', '.bmp', '.tiff']:
            img = Image.open(file_path)
            return pytesseract.image_to_string(img)

        else:
            #text = textract.process(file_path)
            #return text.decode('utf-8')
            text =f"error here file type - {ext} file path - {file_path}"
            return text
            pass

    except Exception as e:
        return f"Error reading file {file_path}: {e}"

def w2v_tokenize_text(text):
    tokens = []
    for sent in nltk.sent_tokenize(text, language='english'):
        for word in nltk.word_tokenize(sent, language='english'):
            if len(word) < 2 or len(word) > 12:
                continue
            tokens.append(word)
    return tokens

def word_averaging(wv, words):
    all_words, mean = set(), []
    for word in words:
        if isinstance(word, np.ndarray):
            mean.append(word)
        elif word in wv.key_to_index:
            mean.append(wv.get_vector(word, norm=True))
            all_words.add(wv.key_to_index[word])
    if not mean:
        return np.zeros(wv.vector_size, dtype=np.float32)
    mean = np.array(mean).mean(axis=0)
    return mean.astype(np.float32)

def word_averaging_list(wv, text_list):
    return np.vstack([word_averaging(wv, post) for post in text_list])

def infer_vector(text, model):
    cleaned_text = clean_text(text)
    tokens = cleaned_text.split()
    return model.infer_vector(tokens)

class WeightedVotingEnsemble(BaseEstimator, ClassifierMixin):
    def __init__(self, models, weights, label_binarizer, labels=None, voting='hard'):
        self.models = models
        self.weights = np.array(weights) / np.sum(weights)
        self.lb = label_binarizer
        self.labels = labels
        self.voting = voting  # 'hard' or 'soft'
    
    def predict(self, X_list):
        if self.voting == 'soft':
            return self._predict_soft(X_list)
        else:
            return self._predict_hard(X_list)

    def _predict_hard(self, X_list):
        all_preds = []
        for model, X in zip(self.models, X_list):
            pred = model.predict(X)
            all_preds.append(pred)
        
        one_hot_preds = np.array([self.lb.transform(p) for p in all_preds])
        weighted_sum = np.tensordot(self.weights, one_hot_preds, axes=(0, 0))
        return self.lb.inverse_transform(weighted_sum)

    def _predict_soft(self, X_list):
        all_probs = []
        for model, X in zip(self.models, X_list):
            if hasattr(model, "predict_proba"):
                probs = model.predict_proba(X)
                all_probs.append(probs)
            else:
                # Skip models without probability prediction
                continue
        
        all_probs = np.array(all_probs)
        soft_votes = np.tensordot(self.weights[:len(all_probs)], all_probs, axes=(0, 0))
        return np.argmax(soft_votes, axis=1)


